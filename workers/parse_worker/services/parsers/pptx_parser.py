import io
import tempfile
from typing import List, Dict, Any
from dataclasses import dataclass

import structlog
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from ..storage import StorageService
from ..parser import ParseResult

logger = structlog.get_logger()

@dataclass
class SlideData:
    slide_number: int
    title: str
    content: str
    notes: str
    image_s3_key: str
    metadata: Dict[str, Any]

class PPTXParser:
    def __init__(self, storage_service: StorageService):
        self.storage = storage_service
    
    async def parse(self, file_data: bytes, deck_id: str) -> ParseResult:
        """Parse PPTX file and extract slides"""
        logger.info("Parsing PPTX file", deck_id=deck_id)
        
        try:
            # Load presentation from bytes
            presentation = Presentation(io.BytesIO(file_data))
            
            slides = []
            metadata = {
                "total_slides": len(presentation.slides),
                "slide_layouts": [],
                "has_notes": False,
                "has_images": False
            }
            
            for i, slide in enumerate(presentation.slides):
                slide_data = await self._extract_slide_data(slide, i + 1, deck_id)
                slides.append(slide_data)
                
                # Update metadata
                if slide_data.notes:
                    metadata["has_notes"] = True
                if slide_data.image_s3_key:
                    metadata["has_images"] = True
                
                # Track slide layouts
                layout_name = slide.slide_layout.name
                if layout_name not in metadata["slide_layouts"]:
                    metadata["slide_layouts"].append(layout_name)
            
            return ParseResult(
                slides_count=len(slides),
                slides=[slide.__dict__ for slide in slides],
                metadata=metadata
            )
            
        except Exception as e:
            logger.error("PPTX parsing failed", deck_id=deck_id, error=str(e))
            raise
    
    async def _extract_slide_data(self, slide, slide_number: int, deck_id: str) -> SlideData:
        """Extract data from a single slide"""
        title = self._extract_title(slide)
        content = self._extract_content(slide)
        notes = self._extract_notes(slide)
        image_s3_key = await self._extract_images(slide, deck_id, slide_number)
        
        metadata = {
            "layout": slide.slide_layout.name,
            "shapes_count": len(slide.shapes),
            "has_images": bool(image_s3_key)
        }
        
        return SlideData(
            slide_number=slide_number,
            title=title,
            content=content,
            notes=notes,
            image_s3_key=image_s3_key,
            metadata=metadata
        )
    
    def _extract_title(self, slide) -> str:
        """Extract title from slide"""
        title = ""
        
        # Try to find title placeholder
        for shape in slide.shapes:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # Title
                if hasattr(shape, "text"):
                    title = shape.text.strip()
                    break
        
        # If no title placeholder, try to find text in first text box
        if not title:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
        
        return title
    
    def _extract_content(self, slide) -> str:
        """Extract content from slide"""
        content_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Skip title
                if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:
                    continue
                
                content_parts.append(shape.text.strip())
        
        return "\n".join(content_parts)
    
    def _extract_notes(self, slide) -> str:
        """Extract speaker notes from slide"""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text.strip()
        return ""
    
    async def _extract_images(self, slide, deck_id: str, slide_number: int) -> str:
        """Extract images from slide and upload to S3"""
        images = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # Get image data
                    image = shape.image
                    image_bytes = image.blob
                    
                    # Convert to PIL Image for processing
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    
                    # Save image to S3
                    image_key = f"decks/{deck_id}/slides/{slide_number}/images/{len(images)}.png"
                    await self.storage.upload_file(image_key, image_bytes, "image/png")
                    
                    images.append(image_key)
                    
                except Exception as e:
                    logger.warning("Failed to extract image", error=str(e))
        
        return images[0] if images else ""
